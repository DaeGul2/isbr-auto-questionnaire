import os
import json
import pandas as pd
from flask import Flask, request, send_file, jsonify
from flask_cors import CORS  # ✅ CORS 라이브러리 추가
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ✅ Flask 서버 초기화
app = Flask(__name__)
CORS(app)

# ✅ 업로드 및 출력 디렉토리 설정
UPLOAD_FOLDER = "uploads"
OUTPUT_FOLDER = "outputs"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER
app.config["OUTPUT_FOLDER"] = OUTPUT_FOLDER

#
# 첫 번째 인자 (Cm(1.3)): 텍스트 박스의 왼쪽 상단 모서리의 x 좌표입니다. 이 값은 슬라이드의 왼쪽 가장자리에서 텍스트 박스 시작점까지의 가로 거리를 센티미터 단위로 나타냅니다.

# 두 번째 인자 (Cm(0.94)): 텍스트 박스의 왼쪽 상단 모서리의 y 좌표입니다. 이 값은 슬라이드의 상단 가장자리에서 텍스트 박스 시작점까지의 세로 거리를 센티미터 단위로 나타냅니다.

# 세 번째 인자 (Cm(7.62)): 텍스트 박스의 너비입니다. 이 값은 텍스트 박스의 가로 길이를 센티미터 단위로 설정합니다.

# 네 번째 인자 (Cm(0.77)): 텍스트 박스의 높이입니다. 이 값은 텍스트 박스의 세로 길이를 센티미터 단위로 설정합니다.
#

# ✅ PPT 생성 함수 (질문 포함)
def create_ppt(excel_path, output_pptx_path, id_to_text):
    df = pd.read_excel(excel_path)
    print(df)
    prs = Presentation()
    # A4 크기 세로 설정 (210mm x 297mm)
    prs.slide_width = Cm(21.0)
    prs.slide_height = Cm(29.7)
    print("컬럼 확인 : ", df.columns)

    whole_width = Cm(19.05)  # 전체 너비를 이 변수로 제어

    for _, row in df.iterrows():
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # 백지 슬라이드

        # ✅ 지원자 ID 삽입 (좌측 상단)
        # 텍스트 박스 설정
        # ✅ 지원자 ID 삽입 (좌측 상단) - 위치 수정
        
        text_box = slide.shapes.add_textbox(Cm(0.98), Cm(0.14), Cm(7.62), Cm(0.77))
        text_frame = text_box.text_frame
        text_frame.clear()

        # 단락 생성
        p = text_frame.add_paragraph()
        p.level = 0  # 들여쓰기 레벨 설정 (0은 들여쓰기 없음)
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(0)  # 문단 앞 공간 제거
        p.space_after = Pt(0)   # 문단 뒤 공간 제거
        

        # '수험번호:' 부분 추가 (볼드 처리)
        run = p.add_run()
        run.text = "면접번호 : "
        run.font.bold = True
        run.font.size = Pt(10)
        run.font.name = "맑은 고딕"

        # 수험번호 값 추가 (볼드 미처리)
        run = p.add_run()
        run.text = f"{row['지원자_ID']}\n\n"
        run.font.bold = False
        run.font.size = Pt(10)
        run.font.name = "맑은 고딕"

        

        # '지원분야:' 부분 추가 (볼드 처리)
        run = p.add_run()
        run.text = "지원분야 : "
        run.font.bold = True
        run.font.size = Pt(10)
        run.font.name = "맑은 고딕"

        # 지원분야 값 추가 (볼드 미처리)
        run = p.add_run()
        run.text = f"{dict(row).get('지원분야', '기본값')}"
        run.font.bold = False
        run.font.size = Pt(10)
        run.font.name = "맑은 고딕"

        # 자기소개서 제목
        title_y = Cm(2.54)
        text_box = slide.shapes.add_textbox(Cm(0.98), Cm(2.21), whole_width, Cm(1.27))
        text_frame = text_box.text_frame
        text_frame.clear()

        p = text_frame.add_paragraph()
        p.text = "자기소개서"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.font.name = "맑은 고딕"
        p.alignment = PP_ALIGN.LEFT
    
        # 자소서_ID에 맞는 설명
        description_y = title_y + Cm(1.27)
        description_text = id_to_text.get(str(row["자소서_ID"]), "기본값")

        text_box = slide.shapes.add_textbox(Cm(0.98), Cm(3.41), whole_width, Cm(1.27))
        text_frame = text_box.text_frame
        text_frame.clear()
        text_frame.word_wrap = True

        p = text_frame.add_paragraph()
        p.text = "질문 : " + description_text
        p.font.size = Pt(10)
        p.font.bold = False
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.font.name = "맑은 고딕"
        p.alignment = PP_ALIGN.LEFT

        # Horizon (수평선)
        line_y = title_y + Cm(1.65)
        line_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Cm(0.98), Cm(3.91), whole_width, Cm(0.05)
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = RGBColor(200, 200, 200)
        line_shape.line.fill.background()

        # 원본 텍스트 및 밑줄 인덱스 처리
        original_text = str(row["원본"]).strip() if pd.notna(row["원본"]) else ""
        underline_ranges = json.loads(row["밑줄 인덱스"]) if pd.notna(row["밑줄 인덱스"]) else []

        text_box = slide.shapes.add_textbox(Cm(0.98), Cm(4.05), whole_width, Cm(5.08))
        text_frame = text_box.text_frame
        text_frame.clear()
        text_frame.word_wrap = True

        p = text_frame.add_paragraph()
        p.line_spacing = 1.5
        p.space_after = Pt(10)
        p.space_before = Pt(10)

        last_index = 0
        for underline in underline_ranges:
            start, end = underline["start"], underline["end"]
            if last_index < start:
                run = p.add_run()
                run.text = original_text[last_index:start]

            run = p.add_run()
            run.text = original_text[start:end + 1]
            run.font.underline = True
            run.font.bold = True

            last_index = end + 1

        if last_index < len(original_text):
            run = p.add_run()
            run.text = original_text[last_index:]

        for run in p.runs:
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(0, 0, 0)
            run.font.name = "맑은 고딕"

        # 질문 목록
        questions = str(row["질문"]).strip() if pd.notna(row["질문"]) else "질문 없음"

        text_box = slide.shapes.add_textbox(Cm(0.8), Cm(20.92), whole_width, Cm(1.27))
        text_frame = text_box.text_frame
        text_frame.clear()

        p = text_frame.add_paragraph()
        p.text = "자기소개서 기반 면접 질문 제안"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.font.name = "맑은 고딕"
        p.alignment = PP_ALIGN.LEFT

        line_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Cm(0.98), Cm(22.57), whole_width, Cm(0.05)
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = RGBColor(200, 200, 200)
        line_shape.line.fill.background()

        question_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Cm(0.98), Cm(23), whole_width, Cm(5.8)
        )
        question_box.fill.solid()
        question_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
        question_box.line.color.rgb = RGBColor(180, 180, 180)
        question_box.shadow.inherit = True

        # 바로 아래의 코드가 포인트입니다.
        question_tf = question_box.text_frame
        question_tf.clear()  # clear 해도 빈 단락 하나는 반드시 남아있음.

        # ★★ 여기서 새 단락을 만들지 않고 기존 기본단락을 이용!
        p = question_tf.paragraphs[0]  
        p.line_spacing = 1.5
        p.text = questions
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_after = Pt(10)
        p.space_before = Pt(10)
        p.alignment = PP_ALIGN.LEFT

    prs.save(output_pptx_path)

# ✅ API 엔드포인트 - 파일 업로드 및 질문 목록 전달
@app.route("/generate_ppt", methods=["POST"])
def generate_ppt():
    if "file" not in request.files or "id_to_text" not in request.form:
        return jsonify({"error": "파일 또는 질문 정보가 없습니다."}), 400

    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "파일명이 없습니다."}), 400

    filename = secure_filename(file.filename)
    file_path = os.path.join(app.config["UPLOAD_FOLDER"], filename)
    file.save(file_path)

    id_to_text = json.loads(request.form["id_to_text"])  # ✅ 질문 데이터 파싱

    # ✅ PPT 생성
    output_pptx_path = os.path.join(app.config["OUTPUT_FOLDER"], "output.pptx")
    create_ppt(file_path, output_pptx_path, id_to_text)

    return jsonify({"ppt_url": "http://localhost:5000/download_ppt"}), 200


# ✅ PPT 파일 다운로드 엔드포인트 추가
@app.route("/download_ppt", methods=["GET"])
def download_ppt():
    ppt_path = os.path.join(app.config["OUTPUT_FOLDER"], "output.pptx")

    if not os.path.exists(ppt_path):
        return jsonify({"error": "PPT 파일을 찾을 수 없습니다."}), 404

    return send_file(ppt_path, as_attachment=True, download_name="output.pptx")


# ✅ 서버 실행
if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000, debug=True)
    print("🔥 Flask 서버가 5000번 포트에서 실행 중!")
