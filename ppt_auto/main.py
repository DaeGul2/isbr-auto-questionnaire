import json
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ✅ 엑셀 파일 경로
excel_path = "./data.xlsx"
output_pptx_path = "./test.pptx"

# ✅ 엑셀 데이터 불러오기
df = pd.read_excel(excel_path)

# ✅ PowerPoint 프레젠테이션 생성
prs = Presentation()

# ✅ 슬라이드별 데이터 처리
for _, row in df.iterrows():
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # 백지 슬라이드

    # ✅ 지원자 ID 삽입 (좌측 상단)
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(3), Inches(0.5))
    text_frame = text_box.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = f"지원자 ID: {row['지원자_ID']}"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.LEFT

    # ✅ '자기소개서' 소제목 (2칸 위로 올림)
    title_y = 1.0  # ✅ 기존보다 더 위로 올려서 Horizon과 절대 안 겹침
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(title_y), Inches(6), Inches(0.5))
    text_frame = text_box.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = "자기소개서"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.LEFT

    # ✅ 자기소개서 아래 Horizon (수평선)
    line_y = title_y + 0.65  # ✅ 충분한 간격 추가
    line_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(line_y), Inches(9), Inches(0.02)
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = RGBColor(200, 200, 200)
    line_shape.line.fill.background()

    # ✅ 원본 텍스트 가져오기
    original_text = str(row["원본"]).strip() if pd.notna(row["원본"]) else ""

    # ✅ 밑줄 인덱스 처리
    underline_ranges = json.loads(row["밑줄 인덱스"]) if pd.notna(row["밑줄 인덱스"]) else []

    # ✅ 원본 텍스트 박스
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(line_y-0.1 ), Inches(9), Inches(2))
    text_frame = text_box.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # ✅ 원본 텍스트 추가 (밑줄 + 볼드 처리 + 줄간격 1.5 적용)
    last_index = 0
    paragraph = text_frame.add_paragraph()
    paragraph.line_spacing = 1.5
    paragraph.space_after = Pt(10)  # ✅ 줄간격 1.5 적용
    paragraph.space_before = Pt(10)

    for underline in underline_ranges:
        start, end = underline["start"], underline["end"]
        if last_index < start:
            run = paragraph.add_run()
            run.text = original_text[last_index:start]

        run = paragraph.add_run()
        run.text = original_text[start:end + 1]
        run.font.underline = True
        run.font.bold = True

        last_index = end + 1

    if last_index < len(original_text):
        run = paragraph.add_run()
        run.text = original_text[last_index:]

    for run in paragraph.runs:
        run.font.size = Pt(12)
        run.font.color.rgb = RGBColor(0, 0, 0)
        run.font.name = "맑은 고딕"

    text_height = len(original_text) / 80 * 0.4
    text_box.height = Inches(text_height)

    # ✅ 질문 리스트 가져오기
    questions = str(row["질문"]).strip() if pd.notna(row["질문"]) else "질문 없음"

    # ✅ 원본 박스 아래에 질문목록 배치
    question_y_position = line_y + 0.7 + text_height  # ✅ 원본 박스 아래 충분한 간격 확보

    # ✅ '자기소개서 기반 면접 질문 제안' 소제목 추가 (16pt)
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(question_y_position), Inches(6), Inches(0.5))
    text_frame = text_box.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = "자기소개서 기반 면접 질문 제안"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.LEFT

    # ✅ 자기소개서 기반 면접 질문 제안 아래 Horizon
    line_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(question_y_position + 0.65), Inches(9), Inches(0.02)
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = RGBColor(200, 200, 200)
    line_shape.line.fill.background()

    # ✅ 질문 리스트 박스 (연한 회색 배경 + 테두리 + 그림자 추가)
    question_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(question_y_position + 1), Inches(9), Inches(1.5)
    )
    question_box.fill.solid()
    question_box.fill.fore_color.rgb = RGBColor(245, 245, 245)  # ✅ 연한 회색 배경
    question_box.line.color.rgb = RGBColor(180, 180, 180)  # ✅ 회색 테두리
    question_box.shadow.inherit = True  # ✅ 그림자 추가

    # ✅ 질문 리스트 텍스트 삽입 (볼드 처리 + 줄간격 1.5 적용)
    text_frame = question_box.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    paragraph = text_frame.add_paragraph()
    paragraph.line_spacing = 1.5
    paragraph.text = questions
    paragraph.font.size = Pt(12)
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(0, 0, 0)
    paragraph.font.name = "맑은 고딕"
    paragraph.space_after = Pt(10)  # ✅ 줄간격 1.5 적용
    paragraph.space_before = Pt(10)

# ✅ PPT 저장
prs.save(output_pptx_path)
print(f"✅ PPT 파일이 생성되었습니다: {output_pptx_path}")
